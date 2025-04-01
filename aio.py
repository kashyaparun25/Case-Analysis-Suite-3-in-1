__import__('pysqlite3')
import sys
sys.modules['sqlite3'] = sys.modules.pop('pysqlite3')

import streamlit as st
import os
import re
import tempfile
import logging
import time
import base64
import json
from typing import List, Dict, Any, Union, Optional, Tuple
from dotenv import load_dotenv
import io

# Document processing libraries
from docx import Document
from docx.shared import Inches, Pt
from docx.enum.text import WD_ALIGN_PARAGRAPH
from docx.enum.style import WD_STYLE_TYPE
import PyPDF2
from pptx import Presentation
import docx2txt
from docx.oxml import OxmlElement
from docx.oxml.ns import qn

# CrewAI imports
from crewai import Agent, Task, Crew, Process, LLM
from crewai.tasks import TaskOutput
from crewai.tools import BaseTool, tool
from pydantic import BaseModel, Field
import litellm
from langchain.tools import Tool

# Configure logging
logging.basicConfig(level=logging.INFO)
logger = logging.getLogger(__name__)

# Set page configuration
st.set_page_config(
    page_title="Document Analysis Suite",
    page_icon="📚",
    layout="wide"
)

# Initialize session state variables
if 'active_tab' not in st.session_state:
    st.session_state.active_tab = "Document Analysis"
if 'uploaded_files' not in st.session_state:
    st.session_state.uploaded_files = None
if 'extracted_text' not in st.session_state:
    st.session_state.extracted_text = ""
if 'first_file_text' not in st.session_state:
    st.session_state.first_file_text = ""
if 'breakdown_generated' not in st.session_state:
    st.session_state.breakdown_generated = False
if 'teaching_plan_generated' not in st.session_state:
    st.session_state.teaching_plan_generated = False
if 'board_plan_generated' not in st.session_state:
    st.session_state.board_plan_generated = False
if 'content_type' not in st.session_state:
    st.session_state.content_type = {"type": "Unknown", "confidence": "Low", "indicators": []}

# Load environment variables
#load_dotenv()

# Set the app title based on content type
def get_app_title(content_type):
    if content_type == "Case Study":
        return "📚 Case Study Analysis Suite"
    elif content_type == "Scientific Article":
        return "🔬 Scientific Article Analysis Suite"
    elif content_type == "News Item":
        return "📰 News Item Analysis Suite"
    else:
        return "📄 Document Analysis Suite"

# Page title and description
st.title(get_app_title(st.session_state.content_type["type"]))
st.subheader("Generate comprehensive Teaching Notes, Teaching Plans, and Discussion Frameworks")
st.write("Developed for BIA 568 (Business Intelligence and Analytics) -- Management of A.I. at Stevens Institute of Technology")

st.write("---")

# Sidebar for API key configuration
with st.sidebar:
    st.title("⚙️ Configuration")
    
    api_key = st.text_input("Enter your Gemini API Key", type="password", 
                          help="Required for the AI model to function")
    if api_key:
        os.environ["GEMINI_API_KEY"] = api_key
        os.environ["GOOGLE_API_KEY"] = api_key
        os.environ["LITELLM_MODEL_DEFAULT_PROVIDER"] = "gemini"

        # Configure litellm
        litellm.set_verbose = True
        litellm_config = {
            "model": "gemini/gemini-2.0-flash",
            "api_key": api_key,
            "provider": "gemini"
        }
    
    st.divider()
    
    # Reset button
    if st.button("🔄 Reset Session", use_container_width=True):
        for key in list(st.session_state.keys()):
            del st.session_state[key]
        st.rerun()

#---------------------------- Utility Functions ----------------------------#

def extract_text_from_pdf(file):
    """Extract text content from PDF file"""
    if isinstance(file, bytes):
        file = io.BytesIO(file)
    pdf_reader = PyPDF2.PdfReader(file)
    text = ""
    for page in pdf_reader.pages:
        text += page.extract_text()
    return text

def extract_text_from_docx(file):
    """Extract text content from DOCX file"""
    if isinstance(file, bytes):
        file = io.BytesIO(file)
    return docx2txt.process(file)

def extract_text_from_pptx(file):
    """Extract text content from PPTX file"""
    if isinstance(file, bytes):
        file = io.BytesIO(file)
    prs = Presentation(file)
    text = ""
    for slide in prs.slides:
        for shape in slide.shapes:
            if hasattr(shape, "text"):
                text += shape.text + "\n"
    return text

def extract_text_from_any_file(file):
    """Extract text based on file type"""
    if file.name.endswith('.pdf'):
        return extract_text_from_pdf(file)
    elif file.name.endswith('.docx'):
        return extract_text_from_docx(file)
    elif file.name.endswith(('.pptx', '.ppt')):
        return extract_text_from_pptx(file)
    else:
        return "Unsupported file format"

def create_download_link(content, filename):
    """Create a download link for text content"""
    b64 = base64.b64encode(content.encode()).decode()
    href = f'<a href="data:text/plain;base64,{b64}" download="{filename}" class="download-button">Download {filename}</a>'
    return href

def create_docx_from_markdown(markdown_content, title, document_type):
    """Convert markdown content to a formatted DOCX file"""
    # Create new document
    doc = Document()
    
    # Add title
    title_paragraph = doc.add_paragraph()
    title_paragraph.alignment = WD_ALIGN_PARAGRAPH.CENTER
    title_run = title_paragraph.add_run(title)
    title_run.font.size = Pt(16)
    title_run.bold = True
    
    # Add document type
    type_paragraph = doc.add_paragraph()
    type_paragraph.alignment = WD_ALIGN_PARAGRAPH.CENTER
    type_run = type_paragraph.add_run(f"Document Type: {document_type}")
    type_run.font.size = Pt(12)
    type_run.italic = True
    
    # Add separator
    doc.add_paragraph("_" * 50)
    
    # Process markdown content
    lines = markdown_content.split('\n')
    current_section = None
    
    for line in lines:
        # Handle headers
        if line.startswith('# '):
            p = doc.add_heading(line[2:], level=1)
            current_section = line[2:]
        elif line.startswith('## '):
            p = doc.add_heading(line[3:], level=2)
            current_section = line[3:]
        elif line.startswith('### '):
            p = doc.add_heading(line[4:], level=3)
        # Handle bullet points
        elif line.startswith('* '):
            p = doc.add_paragraph()
            p.style = 'List Bullet'
            p.add_run(line[2:])
        # Handle numbered lists
        elif re.match(r'^\d+\. ', line):
            p = doc.add_paragraph()
            p.style = 'List Number'
            p.add_run(re.sub(r'^\d+\. ', '', line))
        # Regular paragraph
        elif line.strip():
            p = doc.add_paragraph(line)
    
    # Save to a temporary file
    with tempfile.NamedTemporaryFile(delete=False, suffix='.docx') as tmp_file:
        doc.save(tmp_file.name)
        return tmp_file.name

#---------------------------- Content Type Detector ----------------------------#

class ContentTypeDetector:
    def __init__(self, api_key):
        self.api_key = api_key
        self.llm = LLM(
            model='gemini/gemini-2.0-flash',
            api_key=self.api_key,
            provider="gemini"
        )
    
    def create_detector_agent(self):
        return Agent(
            role="Content Type Analyzer",
            goal="Determine the type of document from its content",
            backstory="""You are an expert in content analysis, capable of 
            identifying different types of documents based on their structure, 
            language, and content. You can distinguish between case studies, 
            scientific articles, news items, and other content types with high accuracy.""",
            llm=self.llm,
            verbose=True
        )
    
    def detect_content_type(self, text):
        """Determine the type of content from the text"""
        detection_task = Task(
            description=f"""
            Analyze the following text and determine whether it is:
            1. A case study
            2. A scientific article
            3. A news item
            4. Other (specify)
            
            Identify key indicators that support your classification, such as:
            - Structure (abstract, methods, results for scientific articles)
            - Language patterns (narrative style for case studies, journalistic style for news)
            - Content elements (company information, research data, current events)
            
            Return the content type and confidence level in the following format:
            
            Content Type: [type]
            Confidence: [high/medium/low]
            Indicators: [list key indicators]
            
            Text to analyze:
            {text[:3000]}
            """,
            expected_output="Content type classification with indicators",
            agent=self.create_detector_agent()
        )
        
        crew = Crew(
            agents=[self.create_detector_agent()],
            tasks=[detection_task],
            process=Process.sequential,
            verbose=False
        )
        result = crew.kickoff()
        
        # Parse the result to extract content type
        content_type = "Case Study"  # Default
        confidence = "Medium"
        indicators = []
        
        for line in result.raw.split('\n'):
            if line.startswith('Content Type:'):
                content_type = line.replace('Content Type:', '').strip()
            elif line.startswith('Confidence:'):
                confidence = line.replace('Confidence:', '').strip()
            elif line.startswith('Indicators:'):
                # Get all indicators (might be on multiple lines)
                indicator_text = result.raw.split('Indicators:')[1].strip()
                indicators = [ind.strip() for ind in indicator_text.split('-') if ind.strip()]
        
        return {
            "type": content_type,
            "confidence": confidence,
            "indicators": indicators
        }

def process_files(uploaded_files):
    """Process uploaded files and extract text with content type detection"""
    combined_text = ""
    first_file_text = ""
    
    temp_file_paths = []
    
    for file in uploaded_files:
        file.seek(0)  # Reset file pointer
        
        # Extract text based on file type
        if file.name.endswith('.pdf'):
            text = extract_text_from_pdf(file)
        elif file.name.endswith('.docx'):
            text = extract_text_from_docx(file)
        elif file.name.endswith(('.pptx', '.ppt')):
            text = extract_text_from_pptx(file)
        else:
            continue
        
        # Save first file's text for metadata extraction
        if not first_file_text:
            first_file_text = text
        
        combined_text += text + "\n\n"
        
        # Create temporary file for each uploaded file
        with tempfile.NamedTemporaryFile(delete=False, suffix=f".{file.name.split('.')[-1]}") as tmp_file:
            file.seek(0)
            tmp_file.write(file.getvalue())
            temp_file_paths.append(tmp_file.name)
    
    # Detect content type if API key is available
    content_type = {"type": "Unknown", "confidence": "Low", "indicators": []}
    if os.environ.get("GEMINI_API_KEY") and first_file_text:
        detector = ContentTypeDetector(os.environ.get("GEMINI_API_KEY"))
        content_type = detector.detect_content_type(first_file_text)
    
    return combined_text, first_file_text, temp_file_paths, content_type

#---------------------------- Document Generator (Document Analysis) ----------------------------#

class DocumentGenerator:
    def __init__(self):
        self.bullet_counter = 1

    def add_toc(self, doc):
        """Add native Word table of contents"""
        paragraph = doc.add_paragraph()
        run = paragraph.add_run()

        # Start the TOC field
        fldChar1 = create_element('w:fldChar')
        create_attribute(fldChar1, 'w:fldCharType', 'begin')
        run._r.append(fldChar1)

        # Add TOC instruction text
        instrText = create_element('w:instrText')
        create_attribute(instrText, 'xml:space', 'preserve')
        instrText.text = 'TOC \\o "1-3" \\h \\z \\u'
        run._r.append(instrText)

        # End the TOC field
        fldChar2 = create_element('w:fldChar')
        create_attribute(fldChar2, 'w:fldCharType', 'end')
        run._r.append(fldChar2)

    def setup_document_styles(self, doc):
        """Set up custom styles for the document"""
        styles = doc.styles
        
        # Heading styles
        for level in range(1, 4):
            style_name = f'Heading {level}'
            if style_name not in styles:
                style = styles.add_style(style_name, WD_STYLE_TYPE.PARAGRAPH)
                style.base_style = styles['Normal']
                style.font.size = Pt(16 - (level * 2))
                style.font.bold = True

        # Custom bullet point style
        if 'Bullet Point' not in styles:
            bullet_style = styles.add_style('Bullet Point', WD_STYLE_TYPE.PARAGRAPH)
            bullet_style.base_style = styles['Normal']
            bullet_style.font.size = Pt(11)
            bullet_style.paragraph_format.left_indent = Inches(0.25)
            bullet_style.paragraph_format.first_line_indent = Inches(-0.25)

    def add_formatted_text(self, paragraph, text):
        """Add text to paragraph with proper formatting"""
        # First handle double asterisks
        parts = re.split(r'(\*\*.*?\*\*)', text)
        
        for part in parts:
            if part.startswith('**') and part.endswith('**'):
                # Handle bold text (surrounded by double asterisks)
                run = paragraph.add_run(part[2:-2])
                run.bold = True
            else:
                # Handle single asterisks within the remaining text
                # Split by single asterisks
                subparts = re.split(r'(\*[^\*]+\*)', part)
                for subpart in subparts:
                    if subpart.startswith('*') and subpart.endswith('*') and len(subpart) > 2:
                        # It's a bold text marked with single asterisks
                        run = paragraph.add_run(subpart[1:-1])
                        run.bold = True
                    else:
                        # Regular text
                        if subpart.strip():
                            paragraph.add_run(subpart)

    def clean_content(self, text):
        """Clean content and formatting"""
        # Remove HTML tags
        text = re.sub(r'<[^>]+>', '', text)
        
        # Remove markdown headers while preserving content
        text = re.sub(r'^#+\s*(.+)$', r'\1', text, flags=re.MULTILINE)
        
        # Remove duplicate section headers
        text = re.sub(r'(?i)^(.*?)\n\*\*\1\*\*', r'\1', text, flags=re.MULTILINE)
        
        # Clean up <br> tags
        text = re.sub(r'<br\s*/?>', '\n', text)
        
        # Remove excessive newlines
        text = re.sub(r'\n\s*\n', '\n\n', text)
        
        # Process line by line to handle bullet points and bold text
        lines = text.split('\n')
        cleaned_lines = []
        
        for line in lines:
            line = line.strip()
            if not line:
                cleaned_lines.append(line)
                continue
                
            # Check if line starts with a single asterisk (potential bullet point)
            if line.lstrip().startswith('*'):
                # Skip if it starts with double asterisks
                if line.lstrip().startswith('**'):
                    cleaned_lines.append(line)
                    continue
                    
                # Count asterisks that are not part of bold text markers
                # First, temporarily replace bold text markers
                temp_line = re.sub(r'\*\*.*?\*\*', '', line)  # Remove double-asterisk patterns
                temp_line = re.sub(r'\*[^\*]+\*', '', temp_line)  # Remove single-asterisk patterns
                
                # If there's exactly one asterisk left, it's a bullet point
                if temp_line.count('*') == 1:
                    content = line.replace('*', '', 1).strip()
                    cleaned_lines.append(f'* {content}')
                else:
                    cleaned_lines.append(line)
            else:
                cleaned_lines.append(line)
        
        return '\n'.join(cleaned_lines).strip()

    def add_section_content(self, doc, section_name, content):
        """Add a section with proper formatting"""
        # Add section heading
        heading = doc.add_heading(section_name.upper(), level=1)
        heading.alignment = WD_ALIGN_PARAGRAPH.LEFT
        
        # Clean and process the content
        content = self.clean_content(content)
        lines = content.split('\n')
        
        for line in lines:
            line = line.strip()
            if not line:
                continue
            
            # Check for bullet point line
            is_bullet = line.startswith('* ') and not line.startswith('** ')
            
            if is_bullet:
                # Create bullet point paragraph
                bullet_paragraph = doc.add_paragraph(style='Bullet Point')
                bullet_paragraph.paragraph_format.alignment = WD_ALIGN_PARAGRAPH.JUSTIFY
                bullet_paragraph.paragraph_format.left_indent = Inches(0.5)
                bullet_paragraph.paragraph_format.first_line_indent = Inches(-0.25)
                
                # Add bullet character
                bullet_paragraph.add_run('• ')
                
                # Add the rest of the line with formatting
                content = line[2:].strip()
                self.add_formatted_text(bullet_paragraph, content)
            else:
                # Regular paragraph
                p = doc.add_paragraph()
                p.paragraph_format.alignment = WD_ALIGN_PARAGRAPH.JUSTIFY
                self.add_formatted_text(p, line)

    def create_word_document(self, title, author, sections_content, document_type):
        """Create and return a formatted Word document"""
        doc = Document()
        self.setup_document_styles(doc)
        
        # Title section
        title_paragraph = doc.add_paragraph()
        title_paragraph.alignment = WD_ALIGN_PARAGRAPH.CENTER
        
        # Set title based on document type
        if document_type == "Case Study":
            title_run = title_paragraph.add_run("CASE BREAKDOWN")
        elif document_type == "Scientific Article":
            title_run = title_paragraph.add_run("ARTICLE ANALYSIS")
        elif document_type == "News Item":
            title_run = title_paragraph.add_run("NEWS ANALYSIS")
        else:
            title_run = title_paragraph.add_run("DOCUMENT ANALYSIS")
            
        title_run.font.size = Pt(16)
        title_run.bold = True
        
        # Document name and author
        case_name_paragraph = doc.add_paragraph()
        case_name_paragraph.alignment = WD_ALIGN_PARAGRAPH.CENTER
        case_name_run = case_name_paragraph.add_run(title)
        case_name_run.font.size = Pt(14)
        case_name_run.italic = True
        
        author_paragraph = doc.add_paragraph()
        author_paragraph.alignment = WD_ALIGN_PARAGRAPH.CENTER
        author_run = author_paragraph.add_run(author)
        author_run.font.size = Pt(12)
        
        # Separator line
        separator = doc.add_paragraph()
        separator.alignment = WD_ALIGN_PARAGRAPH.CENTER
        separator_run = separator.add_run('_' * 50)
        
        # Add TOC
        toc_heading = doc.add_heading('TABLE OF CONTENTS', level=1)
        toc_heading.alignment = WD_ALIGN_PARAGRAPH.CENTER
        self.add_toc(doc)
        doc.add_page_break()
        
        # Add sections
        for section_name, content in sections_content.items():
            self.add_section_content(doc, section_name, content)
            doc.add_paragraph()  # Add spacing between sections
        
        return doc

# Helper functions for Word document
def create_element(name):
    return OxmlElement(name)

def create_attribute(element, name, value):
    element.set(qn(name), value)

#---------------------------- Document Analysis Generator ----------------------------#

class DocumentMetadata(BaseModel):
    title: str = Field(description="The title of the document")
    author: str = Field(description="The author(s) of the document")

class SectionContent(BaseModel):
    content: str = Field(description="Generated content for the section")
    review: str = Field(description="Review of the content with score and feedback")

# Define section templates based on content type
def get_sections_by_content_type(content_type):
    if content_type == "Case Study":
        return {
            "Case Synopsis": """Summarize the case in 3–4 paragraphs of a total 300 words. Explain the company's background, the industry it operates in, 
                            and the key challenge or strategic decision it faces. Discuss any turning points or dilemmas. Why is this case relevant for 
                            business students? 
                            Position in Course: In 1–2 sentences, describe which type of course this case is best suited for 
                            (e.g., Operations Strategy, AI in Business, Global Supply Chain). What key topics does it help students understand""",

            "Learning Objectives": """List 8–10 learning objectives for this case. A total of 300 words.  What should students understand after analyzing this case? Focus on 
                                leadership decisions, operational insights, AI adoption, or ethical concerns""",
            
            "Teaching Strategies":
            """
            Describe 8–10 ways an instructor can effectively teach this case. Should they use role-playing? A total of 350 words.
            A structured debate? Analyzing real-world examples? How can students actively engage with the material? 1st point should be Approach and second point should be Objective
            """,
            
            "Suggested Teaching Plan":
            """
            Outline a structured teaching plan for this case. What should be covered first? A total of 350 words.
            When should student activities be introduced? How should the case discussion conclude?
            """,

            "Key Points and Insights":
            """
            List 10 key insights students should take from this case. What makes this case unique? A total of 300 words.
            What are the most important strategic, operational, or ethical considerations?
            """,

            "Further Insights":
            """
            Provide additional insights that go beyond the case. How does this case relate to 
            larger business trends? What external factors (e.g., regulation, innovation, geopolitical forces) could impact the situation? A total of 400 words.
            """,

            "Discussion Questions & Answers":
            """
            Create 8–10 discussion questions about the case. What are the key strategic dilemmas?
             Where do trade-offs exist? How can students critically analyze the company's decisions? Provide concise answers. A total of 450 words.
            """,

            "Assignment Exercises":
            """
            List 8–10 in-depth assignments that encourage strategic thinking, data analysis, 
            or real-world application. Include a mix of strategy proposals, financial analysis, 
            role-playing exercises, and ethical debates. Define the expected format 
            (e.g., business report, presentation, comparative analysis) and objective of each exercise. A total of 550- 600 words.
            """,

            "Automated Conversation":
            """
            Write an AI-generated conversation between three personas relevant to the case. A total of 650- 700 words.
            One should be an executive making a strategic decision, another should be an expert,
            and the third should be skeptical or resistant. The conversation should explore the challenges, risks, and strategic importance of the case.
            """,

            "Case Suggestions":
            """
            List 10 ways to make this case more interactive. Should students role-play as executives? 
            Simulate a crisis response? Conduct a competitive market analysis? 
            Suggest unique, hands-on ways to engage with the case. A total of 550 words.
            """
        }
        
    elif content_type == "Scientific Article":
        return {
            "Article Summary": """Summarize the scientific article in 3-4 paragraphs totaling 300 words. Explain the research question, 
                                methodology, key findings, and significance. Discuss the implications for the field.
                                Position in Course: In 1-2 sentences, describe which type of course this article is 
                                best suited for and what key topics it helps students understand.""",
                                
            "Research Methodology Analysis": """Analyze the research methodology in 300 words. What approaches were used? 
                                            What were the strengths and limitations? How could the methodology be improved?
                                            Focus on sample selection, data collection methods, analytical techniques, and validation.""",
                                            
            "Key Findings and Implications": """Summarize the 8-10 most important findings and their implications in 350 words.
                                            What new knowledge does this research contribute? How does it confirm or challenge
                                            existing theories? What are the practical applications of these findings?""",
                                            
            "Learning Objectives": """List 8-10 learning objectives for this scientific article in 300 words.
                                    What should students understand after analyzing this article? Focus on research methods, 
                                    data analysis, scientific reasoning, or ethical implications.""",
                                    
            "Teaching Strategies": """Describe 8-10 ways an instructor can effectively teach this article in 350 words.
                                    Should they use research replication? Critical analysis? Applying findings to case studies?
                                    How can students actively engage with the research methodology and results?""",
                                    
            "Classroom Activities": """Outline 8-10 specific classroom activities in 400 words that help students
                                    engage with the research, methodology, and findings. Include individual and group
                                    activities that develop critical thinking and analytical skills.""",
                                    
            "Discussion Questions & Answers": """Create 8-10 discussion questions about the article in 450 words.
                                            What are the key methodological considerations? How do the findings relate to existing
                                            research? What future research might build on these findings? Provide concise answers.""",
            
            "Simulated Research Discussion":
            """
            Write a simulated research discussion between a principal investigator and 3-4 lab members/colleagues analyzing 
            this scientific article. A total of 500-600 words. Include methodological critiques, alternative interpretations of data, 
            suggestions for follow-up experiments, and connections to related research areas. The discussion should demonstrate 
            scientific thinking, healthy skepticism, and collaborative problem-solving. Include moments of disagreement, 
            clarification, and breakthrough insights that advance understanding of the research. You can also create persons and their roles based on the article.
            """,
                                            
            "Further Research Directions": """Suggest 8-10 ways the research could be extended or applied in 400 words.
                                            What questions remain unanswered? What methodology improvements could be made?
                                            How could the findings be tested in different contexts?""",
                                            
            "Interdisciplinary Connections": """Explore how this research connects to 8-10 other disciplines or fields in 450 words.
                                            How might these findings impact or be applied in other domains?
                                            What cross-disciplinary research opportunities exist?""",
                                            
            "Critical Evaluation": """Provide a critical evaluation of the article in 500 words.
                                    Assess the strength of the evidence, validity of conclusions, and overall contribution
                                    to the field. What are the article's strengths and limitations?"""
        }
        
    elif content_type == "News Item":
        return {
            "News Summary": """Summarize the news item in 3-4 paragraphs totaling 300 words.
                            Explain the key events, stakeholders involved, and broader context.
                            Discuss the significance and implications. In 1-2 sentences, describe why
                            this news item is relevant for students and what course topics it relates to.""",
                            
            "Contextual Analysis": """Provide historical and current context for this news item in 350 words.
                                What events, trends, or policies led to this situation? How does this
                                news fit into broader patterns or developments in this field?""",
                                
            "Key Stakeholders Analysis": """Identify and analyze 8-10 key stakeholders in 400 words.
                                        For each stakeholder, explain their interests, influence, positions,
                                        and how they're affected by the events described in the news item.""",
                                        
            "Impact Assessment": """Assess the immediate and potential long-term impacts in 350 words.
                                What are the economic, social, political, and environmental implications?
                                How might different industries, communities, or policies be affected?""",
                                
            "Learning Objectives": """List 8-10 learning objectives in 300 words. What should students
                                    understand after analyzing this news item? Focus on analytical skills,
                                    media literacy, critical thinking, or connecting theory to current events.""",
                                    
            "Discussion Framework": """Create a structured framework for discussing this news item in 400 words.
                                    How should the discussion be organized? What key questions should guide
                                    the conversation? How can instructors ensure balanced perspectives?""",
                                    
            "Discussion Questions & Answers": """Create 8-10 discussion questions in 450 words.
                                            What critical thinking questions will help students analyze this news?
                                            How can students evaluate different perspectives? Provide concise answers.""",
            
            "Simulated Panel Discussion":
            """
            Write a simulated panel discussion between a moderator and 3-4 experts analyzing this news item. 
            A total of 700-800 words. Include diverse professional and political perspectives on the events, causes, 
            implications, and media coverage. The panelists should represent different stakeholders or viewpoints,
            with the moderator guiding the conversation through key aspects of the news. Include moments of civil 
            disagreement, fact-checking, contextualizing of information, and insights into broader trends related to the news item.
            You can also create persons and their roles based on the article.
            """,
        
                                            
            "Classroom Activities": """Suggest 8-10 classroom activities in 500 words related to this news item.
                                    Include debates, simulations, research projects, media analysis exercises,
                                    and other activities that promote deep engagement with the content.""",
                                    
            "Media Analysis Component": """Analyze how this news has been covered by different sources in 400 words.
                                        Compare coverage across 3-4 different media outlets. Identify potential
                                        biases, framing choices, and what might be emphasized or omitted.""",
                                        
            "Related Resources": """Identify 8-10 related resources in 350 words for further exploration.
                                Include academic articles, books, documentaries, podcasts, and other
                                news sources that provide deeper context or alternative perspectives."""
        }
    else:
        return {
            "Content Summary": """Summarize the document in 3-4 paragraphs totaling 300 words.
                                Explain the key themes, information, and purpose of the content.
                                Discuss why this content is valuable for teaching and learning.""",
                                
            "Key Points and Themes": """Identify 8-10 key points and themes in 350 words.
                                    What are the most important ideas, concepts, or arguments presented?
                                    What makes this content valuable to study?""",
                                    
            "Learning Objectives": """List 8-10 learning objectives in 300 words. What should students
                                    understand after studying this content? What skills might they develop?
                                    How does this content connect to broader educational goals?""",
                                    
            "Teaching Approaches": """Suggest 8-10 teaching approaches in 400 words.
                                    How can this content be effectively taught? What instructional
                                    methods would be most appropriate for this material?""",
                                    
            "Discussion Questions": """Create 8-10 discussion questions in 450 words.
                                    What questions will stimulate critical thinking and deep engagement
                                    with the content? Provide brief answers or guidance for each.""",
                                    
            "Learning Activities": """Propose 8-10 learning activities in 500 words.
                                    What individual and group activities would help students
                                    engage with and apply the content? Include a mix of analytical,
                                    creative, and practical exercises.""",
                                    
            "Further Exploration": """Suggest resources and directions for further exploration in 350 words.
                                    What related topics might students investigate? What additional
                                    readings or materials would complement this content?""",

            "Simulated Dialogue":
            """
            Write a simulated dialogue between an educator and 3-4 participants analyzing this content.
            A total of 700-800 words. Include different interpretations, connections to prior knowledge,
            practical applications, and critical questions about the material. The dialogue should demonstrate
            how different perspectives can enrich understanding and how thoughtful questioning can reveal deeper 
            insights about the content. Include moments of clarification, realization, and connection-making.
            You can also create persons and their roles based on the article.
            """,
                                    
            "Relevance and Applications": """Discuss the relevance and practical applications in 400 words.
                                        How does this content apply to real-world situations?
                                        Why is it important for students to engage with this material?""",
                                        
            "Critical Analysis": """Provide a critical analysis of the content in 450 words.
                                What perspectives are represented or missing? What assumptions
                                underlie the content? What strengths and limitations does it have?""",
                                
            "Assessment Strategies": """Suggest 8-10 assessment strategies in 400 words.
                                    How might instructors evaluate student understanding and
                                    application of this content? Include formal and informal
                                    assessment approaches."""
        }

class DocumentAnalysisCrew:
    def __init__(self, api_key):
        self.api_key = api_key
        self.llm = LLM(
            model='gemini/gemini-2.0-flash',
            api_key=self.api_key,
            provider="gemini"
        )
    
    def create_metadata_agent(self):
            return Agent(
                role="Metadata Analyzer",
                goal="Extract title and author information from document content",
                backstory="""You specialize in analyzing document content to identify key metadata
                such as titles, authors, and other publication information. You have a keen eye for
                identifying the most important and relevant document metadata, even when it's not
                explicitly labeled.""",
                llm=LLM(
                model='gemini/gemini-2.0-flash',
                api_key=self.api_key,
                provider="gemini"
            ),
                verbose=True
            )
        
    def create_content_generator_agent(self):
        
        return Agent(
            role="Document Content Generator",
            goal="Generate comprehensive analysis content based on section requirements",
            backstory="""You are an expert analyst specializing in creating educational content from various
            document types including case studies, scientific articles, and news items. You excel at breaking 
            down complex content into structured, insightful analysis that highlights key learning points, research
            findings, or current events context. You have extensive experience in education and know how to create 
            content that is valuable for teaching and learning.""",
            llm=self.llm,
            verbose=True
        )
    
    def create_content_reviewer_agent(self):
        
        return Agent(
            role="Content Quality Reviewer",
            goal="Evaluate and score content for quality, relevance, and depth",
            backstory="""You are a seasoned academic reviewer with years of experience evaluating
            educational content across various formats. You have a strong understanding of what makes
            effective teaching material for different document types and can provide constructive feedback 
            to improve content quality. You carefully analyze content for relevance, clarity, depth, and 
            educational value.""",
            llm=self.llm,
            verbose=True
        )
    
    def create_metadata_task(self, text):
        return Task(
            description=f"""
            Analyze the extracted text and identify the document title and author(s).
            
            Look for information typically found at the beginning of a document, such as:
            1. The document title (article title, case study name, news headline)
            2. The author name(s)
            
            Return the identified title and author in the following format:
            
            Title: [the title]
            Author: [the author(s)]
            
            If you cannot find this information with certainty, use "Untitled Document" for the title
            and "Unknown Author" for the author.
            
            Text to analyze:
            {text[:2000]}
            """,
            expected_output="Extracted metadata with title and author information",
            agent=self.create_metadata_agent()
        )
    
    def create_section_task(self, section_name, prompt, text, content_type):
        # Adjust prompt based on content type and section
        # Content type adjustments are handled at the section template level now
        
        return Task(
            description=f"""
            Generate content for the '{section_name}' section of the document analysis.
            
            {prompt}
            
            Formatting Requirements:
            1. Use **bold** for important terms or concepts
            2. For bullet points, start each line with "* " (asterisk followed by space)
            3. For numbered lists, use "1. ", "2. " etc.
            4. Use line breaks between paragraphs
            5. Keep paragraphs focused and concise
            
            Remember that this is for a {content_type} document, so tailor your analysis appropriately.
            
            Ensure the content is structured, well-organized, and follows proper formatting.
            
            Text to analyze:
            {text[:5000]}
            """,
            expected_output=f"A well-formatted {section_name} section",
            agent=self.create_content_generator_agent()
        )
    
    def create_review_task(self, section_name, content, content_type):
        return Task(
            description=f"""
            Review the content for the '{section_name}' section for a {content_type} document.
            
            Score it from 1-10 based on:
            1. Relevance to the section (0-3)
            2. Clarity and coherence (0-3)
            3. Depth of analysis (0-4)
            
            Provide a structured review with:
            1. Numerical score
            2. Specific strengths
            3. Areas for improvement
            
            Content to review:
            {content}
            """,
            expected_output=f"A review of the {section_name} section",
            agent=self.create_content_reviewer_agent()
        )
    
    def extract_metadata(self, text):
        metadata_task = self.create_metadata_task(text)
        crew = Crew(
            agents=[self.create_metadata_agent()],
            tasks=[metadata_task],
            process=Process.sequential,
            verbose=False
        )
        result = crew.kickoff()
        
        title = "Untitled Document"
        author = "Unknown Author"
        
        for line in result.raw.split('\n'):
            if line.startswith('Title:'):
                extracted_title = line.replace('Title:', '').strip()
                if extracted_title and extracted_title != "[Untitled Document]":
                    title = extracted_title
            elif line.startswith('Author:'):
                extracted_author = line.replace('Author:', '').strip()
                if extracted_author and extracted_author != "[Unknown Author]":
                    author = extracted_author
        
        return title, author
    
    def generate_section_content(self, text, section_name, section_prompt, content_type):
        section_task = self.create_section_task(section_name, section_prompt, text, content_type)
        crew = Crew(
            agents=[self.create_content_generator_agent()],
            tasks=[section_task],
            process=Process.sequential,
            verbose=False
        )
        result = crew.kickoff()
        return result.raw
    
    def review_content(self, content, section_name, content_type):
        review_task = self.create_review_task(section_name, content, content_type)
        crew = Crew(
            agents=[self.create_content_reviewer_agent()],
            tasks=[review_task],
            process=Process.sequential,
            verbose=False
        )
        result = crew.kickoff()
        return result.raw

#---------------------------- Teaching Plan Generator ----------------------------#

class AgentTracker:
    def __init__(self):
        self.current_agent = ""
        self.placeholder = None
    
    def set_placeholder(self, placeholder):
        self.placeholder = placeholder
    
    def update_agent(self, agent_name):
        self.current_agent = agent_name
        if self.placeholder:
            with self.placeholder:
                st.write(f"🤖 Agent in action: **{self.current_agent}**")

@tool
def extract_text(file_path: str) -> str:
    """
    Extract text from a file based on its extension.
    
    Args:
        file_path (str): Path to the file
        
    Returns:
        str: Extracted text content
    """
    try:
        file_extension = file_path.split('.')[-1].lower()
        
        if file_extension == 'pdf':
            reader = PyPDF2.PdfReader(file_path)
            text = ""
            for page in reader.pages:
                text += page.extract_text()
            return text
        
        elif file_extension == 'docx':
            return docx2txt.process(file_path)
        
        elif file_extension in ['ppt', 'pptx']:
            presentation = Presentation(file_path)
            text = ""
            for slide in presentation.slides:
                for shape in slide.shapes:
                    if hasattr(shape, "text"):
                        text += shape.text + "\n"
            return text
        
        else:
            return f"Unsupported file format: {file_extension}"
    
    except Exception as e:
        return f"Error extracting text: {str(e)}"

def create_teaching_plan_crew(file_paths, content_type):
    # Initialize LLM 
    my_llm = LLM(
        model='gemini/gemini-2.0-flash',
        api_key=os.environ.get("GEMINI_API_KEY"),
        provider="gemini"
    )
    
    # Create tracker instance
    tracker = AgentTracker()
    
    # Adjust agent descriptions and goals based on content type
    if content_type == "Case Study":
        analyzer_role = 'Case Study Analyzer'
        analyzer_goal = 'Extract key concepts, objectives, and data from case study files'
        analyzer_backstory = 'You are an expert in analyzing business case studies. Identify core themes, data points, and learning objectives.'
        
        designer_role = 'Teaching Plan Designer'
        designer_goal = 'Create a 2-hour lesson plan based on analyzed case study data'
        designer_backstory = 'You are an educator designing a structured lesson plan. Use the extracted case study data to outline activities, discussions, and assessments.'
        
        task_description = "Design a 2-hour lesson plan with introduction, case analysis, group activity, and assessment."
    
    elif content_type == "Scientific Article":
        analyzer_role = 'Scientific Article Analyzer'
        analyzer_goal = 'Extract key research findings, methodology, and implications from scientific article'
        analyzer_backstory = 'You are an expert in analyzing scientific research articles. Identify research questions, methodology, findings, and significance.'
        
        designer_role = 'Scientific Article Teaching Designer'
        designer_goal = 'Create a 2-hour research-focused lesson plan based on analyzed article data'
        designer_backstory = 'You are an educator designing a structured lesson plan for scientific content. Use the extracted research data to outline discussions, critiques, and applications.'
        
        task_description = "Design a 2-hour lesson plan with research overview, methodology analysis, findings discussion, and research application activities."
    
    elif content_type == "News Item":
        analyzer_role = 'News Content Analyzer'
        analyzer_goal = 'Extract key events, stakeholders, and context from news item'
        analyzer_backstory = 'You are an expert in analyzing current events and news content. Identify key events, stakeholders, and broader context.'
        
        designer_role = 'Current Events Teaching Designer'
        designer_goal = 'Create a 2-hour current events lesson plan based on analyzed news data'
        designer_backstory = 'You are an educator designing a structured lesson plan for current events. Use the extracted news content to outline discussions, contextual analysis, and media literacy activities.'
        
        task_description = "Design a 2-hour lesson plan with news overview, context discussion, stakeholder analysis, and media literacy activities."
    
    else:
        analyzer_role = 'Content Analyzer'
        analyzer_goal = 'Extract key concepts, themes, and insights from document content'
        analyzer_backstory = 'You are an expert in analyzing various types of documents. Identify key themes, concepts, and learning points.'
        
        designer_role = 'Learning Plan Designer'
        designer_goal = 'Create a 2-hour learning plan based on analyzed document content'
        designer_backstory = 'You are an educator designing a structured learning plan. Use the extracted content to outline discussions, activities, and assessments.'
        
        task_description = "Design a 2-hour lesson plan with content introduction, key concept analysis, interactive activities, and assessment."
    
    # Create agents with Gemini configuration
    pdf_analyzer = Agent(
        role=analyzer_role,
        goal=analyzer_goal,
        backstory=analyzer_backstory,
        llm=my_llm,
        tools=[extract_text],
        verbose=True,
        step_callback=lambda *args, **kwargs: tracker.update_agent(analyzer_role)
    )

    plan_generator = Agent(
        role=designer_role,
        goal=designer_goal,
        backstory=designer_backstory,
        llm=my_llm,
        verbose=True,
        step_callback=lambda *args, **kwargs: tracker.update_agent(designer_role)
    )

    reviewer = Agent(
        role='Plan Reviewer',
        goal='Ensure the lesson plan is engaging and aligned with learning objectives',
        backstory='You are a curriculum reviewer. Verify the plan\'s clarity, alignment with objectives, and engagement level.',
        llm=my_llm,
        verbose=True,
        step_callback=lambda *args, **kwargs: tracker.update_agent("Plan Reviewer")
    )

    final_reporter = Agent(
        role='Teaching Plan Reporter',
        goal='Ensure to incorporate the feedback from the reviewer agent and finalize the content for the teaching plan',
        backstory=f""" You are an expert educator with 20 years of experience in teaching and curriculum development 
                    specializing in {content_type}s. You are responsible for finalizing the content for the teaching plan 
                    and ensuring it is engaging and aligned with learning objectives. You will use the feedback from the 
                    reviewer agent to make necessary revisions and finalize the content for the teaching plan.""",
        llm=my_llm,
        verbose=True,
        step_callback=lambda *args, **kwargs: tracker.update_agent("Teaching Plan Reporter")
    )

    # Combine all file contents into one text
    combined_file_path = file_paths[0]  # Use the first file path for the analyzer to begin
    
    # Define tasks
    analyze_pdf = Task(
        description=f"Extract and analyze the files at {', '.join(file_paths)} for key concepts and learning objectives, keeping in mind this is a {content_type}.",
        config={"file_path": combined_file_path},
        expected_output=f"A summary of key concepts and learning objectives extracted from the {content_type}.",
        agent=pdf_analyzer
    )

    generate_plan = Task(
        description=task_description,
        expected_output=f"A detailed 2-hour lesson plan for a {content_type} with clear sections and activities.",
        agent=plan_generator
    )

    review_plan = Task(
        description=f"Review the lesson plan for clarity, alignment with objectives, and student engagement, considering this is a {content_type}.",
        expected_output="Feedback on the lesson plan's clarity, alignment, and engagement level.",
        agent=reviewer
    )

    final_plan = Task(
        description=f"""Generate the final lesson plan based on the feedback, optimized for teaching a {content_type}.""",
        expected_output="""
                        You are also responsible for ensuring the plan is clear and concise. 
                    - It should have the overall objective to start with.
                    - It should have a clear introduction
                    - It should have detailed lesson breakdowns with the time to be spent on each section and the title of the section Highlighting all the important concepts to be taught (you can use tables to highlight the concepts in a structured way)
                    - It should have one powerful visual aid which can be a table That can help in better understanding for the students
                    - It should include a simulated 10-15 minute class discussion segment showing how a skilled educator might guide students through a key concept in the document
                    - It can have simple assessments that enables class participation engagement in brainstorming activities and such that can help in better understanding of the concepts.
                    - It should have a clear conclusion that ties back to the overall objective.
                    - Overall Plan should be around 1300 - 1500 words.""",
        agent=final_reporter
    )
    
    # Create crew
    crew = Crew(
        agents=[pdf_analyzer, plan_generator, reviewer, final_reporter],
        tasks=[analyze_pdf, generate_plan, review_plan, final_plan],
        process=Process.sequential,
        verbose=True
    )
    
    return crew, tracker

#---------------------------- Discussion Framework Generator ----------------------------#

class DiscussionFrameworkAnalyzer:
    def __init__(self, api_key):
        self.api_key = api_key
        if not api_key:
            raise ValueError("API key not found")
            
        self.llm = LLM(
            model='gemini/gemini-2.0-flash',
            api_key=api_key,
            provider="gemini"
        )
        
        litellm.set_verbose = True

        # Create agents
        self.create_agents()

    def create_agents(self):
        """Create specialized agents for different tasks"""
        
        # Document Processing Agent
        self.document_processor = Agent(
            role='Document Processor',
            goal='Extract and clean text content from document files',
            backstory="""You are an expert at processing various document types and extracting 
            meaningful content. You ensure the text is properly formatted and ready 
            for analysis.""",
            tools=[Tool(
                name="extract_text",
                func=self.extract_text_from_file,
                description="Extracts text content from files"
            )],
            allow_delegation=False,
            llm=self.llm,
            verbose=True
        )

        # Content Analysis Agent
        self.analyzer = Agent(
            role='Document Analyzer',
            goal='Analyze documents and identify key points for discussion framework',
            backstory="""You are an expert analyst skilled at analyzing various types of documents
            and identifying crucial elements for educational discussion. You create clear, structured 
            analyses that highlight key insights for teaching.""",
            tools=[Tool(
                name="analyze_content",
                func=self.analyze_document,
                description="Analyzes document and creates structured discussion framework"
            )],
            allow_delegation=False,
            llm=self.llm,
            verbose=True
        )

    def extract_text_from_file(self, file_content: bytes) -> str:
        """Extract text from file using appropriate method"""
        try:
            # Create file reader object
            file_obj = io.BytesIO(file_content)
            
            # Determine file type
            # This is simplified - in reality, you'd check the file signature/magic bytes
            if b"%PDF" in file_content[:1024]:  # Check for PDF signature
                pdf_reader = PyPDF2.PdfReader(file_obj)
                text_content = []
                for page in pdf_reader.pages:
                    text_content.append(page.extract_text())
                return "\n\n".join(text_content)
            else:
                # Try DOCX
                try:
                    return docx2txt.process(file_obj)
                except:
                    # Try PPTX
                    try:
                        file_obj.seek(0)
                        prs = Presentation(file_obj)
                        text_content = []
                        for slide in prs.slides:
                            for shape in slide.shapes:
                                if hasattr(shape, "text"):
                                    text_content.append(shape.text)
                        return "\n\n".join(text_content)
                    except:
                        return "Unable to extract text from unsupported file format"
            
        except Exception as e:
            raise Exception(f"Extraction error: {str(e)}")

    def analyze_document(self, text: str, document_type: str) -> dict:
        """Analyze document content using litellm with JSON response"""
        # Adjust prompt based on document type
        if document_type == "Case Study":
            framework_type = "BOARD PLAN"
            sections_prompt = """
                Create a detailed analysis with these exact sections:
                1. Main business challenge/opportunity
                2. Industry context and competitors
                3. Key decision points
                4. Stakeholder perspectives
                5. Implementation considerations
                - Resource requirements
                - Timeline
                - Success metrics
                6. Alternative approaches
                7. Learning takeaways
            """
        elif document_type == "Scientific Article":
            framework_type = "DISCUSSION POINTS"
            sections_prompt = """
                Create a detailed analysis with these exact sections:
                1. Main research question/hypothesis
                2. Methodology assessment
                3. Key findings summary
                4. Limitations and validity
                5. Theoretical implications
                6. Practical applications
                7. Future research directions
                8. Interdisciplinary connections
            """
        elif document_type == "News Item":
            framework_type = "CURRENT EVENTS FRAMEWORK"
            sections_prompt = """
                Create a detailed analysis with these exact sections:
                1. Main event/development
                2. Historical and political context
                3. Key stakeholders and perspectives
                4. Media coverage analysis
                5. Immediate impacts
                6. Long-term implications
                7. Related developments
                8. Critical thinking questions
            """
        else:
            framework_type = "DISCUSSION FRAMEWORK"
            sections_prompt = """
                Create a detailed analysis with these exact sections:
                1. Main topic/theme
                2. Key concepts
                3. Important relationships
                4. Supporting evidence
                5. Counterarguments/limitations
                6. Practical applications
                7. Discussion questions
                8. Further exploration topics
            """
        
        messages = [
            {
                "role": "user",
                "content": f"""Analyze this {document_type.lower()} and create a structured {framework_type}.
                
                Content: {text[:5000]}
                
                {sections_prompt}
                
                Also include a "Questions for Further Exploration" section.
                
                Format the response as a JSON object with this structure:
                {{
                    "sections": [
                        {{
                            "title": "SECTION 1: [Title]",
                            "points": ["point 1", "point 2", "point 3"]
                        }},
                        // ... other sections
                    ],
                    "questions": {{
                    "title": "QUESTIONS FOR FURTHER EXPLORATION",
                    "points": ["question 1", "question 2"]
                    }},
                    "discussion": {{
                "title": "SIMULATED DISCUSSION",
                "content": "A 300-400 word simulated discussion between an educator and participants"
            }}
                }}"""
            }
        ]
        
        try:
            response = litellm.completion(
                model="gemini/gemini-2.0-flash",
                messages=messages,
                api_key=os.environ.get("GEMINI_API_KEY"),
                provider="gemini",
                response_format={"type": "json_object"}
            )
            
            # Extract and parse the JSON response
            content = response.choices[0].message.content
            return json.loads(content)
            
        except Exception as e:
            st.error(f"Analysis error: {str(e)}")
            if 'response' in locals():
                st.error("Raw response:")
                st.code(response.choices[0].message.content)
            raise

    def process_document(self, file_content: bytes, document_type: str) -> dict:
        """Process the document using Crew AI agents"""
        try:
            # First extract text
            progress_text = st.empty()
            progress_text.text(f"Extracting text from {document_type.lower()}...")
            text_content = self.extract_text_from_file(file_content)
            
            # Show extracted text for verification
            with st.expander("View extracted text"):
                st.text(text_content[:500] + "...")
            
            # Then analyze content
            progress_text.text(f"Analyzing {document_type.lower()} content...")
            result = self.analyze_document(text_content, document_type)
            
            return result
            
        except Exception as e:
            raise Exception(f"Processing error: {str(e)}")

#---------------------------- Main App Interface ----------------------------#

# File upload section - shared across all generators
if st.session_state.uploaded_files is None:
    uploaded_files = st.file_uploader(
        "Upload document files (PDF, DOCX, PPT, PPTX)", 
        accept_multiple_files=True,
        type=['pdf', 'docx', 'ppt', 'pptx']
    )
    
    if uploaded_files:
        st.session_state.uploaded_files = uploaded_files
        
        # Process files only once and store results
        with st.spinner("Processing uploaded files..."):
            combined_text, first_file_text, temp_file_paths, content_type = process_files(uploaded_files)
            st.session_state.combined_text = combined_text
            st.session_state.first_file_text = first_file_text
            st.session_state.temp_file_paths = temp_file_paths
            st.session_state.content_type = content_type
            
            # Show file upload summary
            st.success(f"✅ {len(uploaded_files)} file(s) uploaded and processed successfully")
            
            # Show uploaded files
            for file in uploaded_files:
                st.write(f"- {file.name}")

# Tabs for different generators with updated names
if st.session_state.uploaded_files:
    # Show detected content type and allow user to change it
    content_type_col, confidence_col = st.columns(2)
    with content_type_col:
        st.info(f"📄 Detected content type: **{st.session_state.content_type['type']}**")
    with confidence_col:
        st.info(f"🔍 Detection confidence: **{st.session_state.content_type['confidence']}**")
    
    # Always display the document type selection
    content_options = ["Case Study", "Scientific Article", "News Item", "Other"]
    selected_type = st.selectbox(
        "Select document type:", 
        options=content_options,
        index=content_options.index(st.session_state.content_type['type']) if st.session_state.content_type['type'] in content_options else 0,
        key="document_type_selector"
    )
    
    # Update the content type in session state
    if selected_type != st.session_state.content_type['type']:
        st.session_state.content_type['type'] = selected_type
        # Reset generated flags when document type changes
        st.session_state.breakdown_generated = False
        st.session_state.teaching_plan_generated = False
        st.session_state.board_plan_generated = False
        st.rerun()
    
    # Update title based on content type
    st.title(get_app_title(st.session_state.content_type["type"]))
    
    # Customize tab names based on content type
    if st.session_state.content_type["type"] == "Case Study":
        tab_names = ["Case Breakdown Generator", "Teaching Plan Generator", "Board Plan Generator"]
    elif st.session_state.content_type["type"] == "Scientific Article":
        tab_names = ["Article Analysis Generator", "Teaching Plan Generator", "Discussion Points Generator"]
    elif st.session_state.content_type["type"] == "News Item":
        tab_names = ["News Analysis Generator", "Discussion Guide Generator", "Current Events Framework"]
    else:
        tab_names = ["Content Analysis Generator", "Teaching Guide Generator", "Discussion Framework"]
    
    # Use native streamlit tabs with content-specific names
    tab1, tab2, tab3 = st.tabs(tab_names)
    
    #------------------ Document Analysis Generator Tab ------------------#
    with tab1:
        st.header(tab_names[0])
        st.write(f"Generate a comprehensive analysis of the {st.session_state.content_type['type'].lower()} with sections for teaching purposes.")
        
        if not st.session_state.breakdown_generated:
            if not api_key:
                st.warning("⚠️ Please enter an API key in the sidebar before proceeding.")
            else:
                # Initialize the breakdown generator
                crew_manager = DocumentAnalysisCrew(api_key)
                
                # Extract metadata
                with st.spinner("Extracting document metadata..."):
                    title, author = crew_manager.extract_metadata(st.session_state.first_file_text)

                # Display extracted metadata with option to edit
                st.subheader("Extracted Document Information")
                title = st.text_input("Document Title:", title)
                author = st.text_input("Author:", author)
                
                # Get sections based on content type
                sections = get_sections_by_content_type(st.session_state.content_type["type"])
                
                # Generate button
                if st.button("Generate Analysis", key="breakdown_button"):
                    # Create tabs for content and review
                    content_tab, review_tab = st.tabs(["Generated Content", "Content Review"])
                    
                    # Generate and display content
                    sections_content = {}
                    reviews = {}
                    
                    with st.spinner("Generating content... This may take several minutes."):
                        progress_bar = st.progress(0)
                        
                        with content_tab:
                            for i, (section_name, prompt) in enumerate(sections.items()):
                                # Generate content using the CrewAI agents
                                st.info(f"Generating {section_name}...")
                                content = crew_manager.generate_section_content(
                                    st.session_state.combined_text, 
                                    section_name, 
                                    prompt, 
                                    st.session_state.content_type["type"]
                                )
                                sections_content[section_name] = content
                                
                                # Generate review
                                review = crew_manager.review_content(
                                    content, 
                                    section_name, 
                                    st.session_state.content_type["type"]
                                )
                                reviews[section_name] = review
                                
                                # Display content in preview
                                st.subheader(section_name)
                                st.markdown(content)
                                
                                # Update progress
                                progress_bar.progress((i + 1) / len(sections))
                    
                    # Store generated content in session state
                    st.session_state.sections_content = sections_content
                    st.session_state.reviews = reviews
                    st.session_state.title = title
                    st.session_state.author = author
                    st.session_state.breakdown_generated = True
                    
                    # Display reviews in review tab
                    with review_tab:
                        for section_name, review in reviews.items():
                            st.subheader(f"{section_name} Review")
                            st.markdown(review)
                    
                    # Generate and store the document for later
                    doc_generator = DocumentGenerator()
                    doc = doc_generator.create_word_document(
                        title=title,
                        author=author,
                        sections_content=sections_content,
                        document_type=st.session_state.content_type["type"]
                    )
                    
                    # Save document to temporary file
                    with tempfile.NamedTemporaryFile(delete=False, suffix='.docx') as tmp_file:
                        doc.save(tmp_file.name)
                        st.session_state.doc_path = tmp_file.name
                        
                    # Allow the user to download the document
                    with open(st.session_state.doc_path, 'rb') as file:
                        # Generate filename based on content type
                        if st.session_state.content_type["type"] == "Case Study":
                            filename = f"{title.lower().replace(' ', '_')}_case_breakdown.docx"
                        elif st.session_state.content_type["type"] == "Scientific Article":
                            filename = f"{title.lower().replace(' ', '_')}_article_analysis.docx"
                        elif st.session_state.content_type["type"] == "News Item":
                            filename = f"{title.lower().replace(' ', '_')}_news_analysis.docx"
                        else:
                            filename = f"{title.lower().replace(' ', '_')}_content_analysis.docx"
                            
                        st.download_button(
                            label=f"📥 Download {st.session_state.content_type['type']} Analysis (DOCX)",
                            data=file,
                            file_name=filename,
                            mime="application/vnd.openxmlformats-officedocument.wordprocessingml.document",
                            key="download_docx_button"
                        )
                
        else:
            # If we already have generated content, display it without regenerating
            st.subheader("Extracted Document Information")
            st.text_input("Document Title:", value=st.session_state.title, key="title_display", disabled=True)
            st.text_input("Author:", value=st.session_state.author, key="author_display", disabled=True)
            
            # Display the generated content in tabs
            content_tab, review_tab = st.tabs(["Generated Content", "Content Review"])
            
            with content_tab:
                for section_name, content in st.session_state.sections_content.items():
                    st.subheader(section_name)
                    st.markdown(content)
            
            with review_tab:
                for section_name, review in st.session_state.reviews.items():
                    st.subheader(f"{section_name} Review")
                    st.markdown(review)
            
            # Allow the user to download the document without regenerating
            with open(st.session_state.doc_path, 'rb') as file:
                # Generate filename based on content type

                if st.session_state.content_type["type"] == "Case Study":
                    filename = f"{st.session_state.title.lower().replace(' ', '_')}_case_breakdown.docx"
                elif st.session_state.content_type["type"] == "Scientific Article":
                    filename = f"{st.session_state.title.lower().replace(' ', '_')}_article_analysis.docx"
                elif st.session_state.content_type["type"] == "News Item":
                    filename = f"{st.session_state.title.lower().replace(' ', '_')}_news_analysis.docx"
                else:
                    filename = f"{st.session_state.title.lower().replace(' ', '_')}_content_analysis.docx"
                    
                st.download_button(
                    label=f"📥 Download {st.session_state.content_type['type']} Analysis (DOCX)",
                    data=file,
                    file_name=filename,
                    mime="application/vnd.openxmlformats-officedocument.wordprocessingml.document",
                    key="download_docx_button" 
                )
    
    #------------------ Teaching Plan Generator Tab ------------------#
    with tab2:
        st.header(tab_names[1])
        st.write(f"Generate a comprehensive teaching plan for the {st.session_state.content_type['type'].lower()}.")
        
        if not st.session_state.teaching_plan_generated:
            if not api_key:
                st.warning("⚠️ Please enter an API key in the sidebar before proceeding.")
            else:
                # Create a button to start generation
                if st.button("Generate Teaching Plan", key="teaching_plan_button"):
                    try:
                        # Create placeholders for UI updates
                        progress_placeholder = st.empty()
                        agent_status_placeholder = st.empty()
                        
                        # Initialize progress bar
                        progress_bar = progress_placeholder.progress(0)
                        
                        # Update progress
                        progress_bar.progress(10)
                        st.info("🔍 Initializing crew and analyzing documents...")
                        
                        # Create crew with content type awareness
                        crew, tracker = create_teaching_plan_crew(
                            st.session_state.temp_file_paths, 
                            st.session_state.content_type["type"]
                        )
                        
                        # Set the tracker's placeholder
                        tracker.set_placeholder(agent_status_placeholder)
                        
                        # Update progress
                        progress_bar.progress(30)
                        st.info("📖 Document analysis in progress...")
                        
                        # Start the execution
                        start_time = time.time()
                        result = crew.kickoff(inputs={"file_path": st.session_state.temp_file_paths[0]})
                        result_text = str(result)  # Convert CrewOutput to string
                        
                        # Update progress
                        progress_bar.progress(90)
                        st.info("✏️ Finalizing teaching plan...")
                        
                        # Calculate execution time
                        execution_time = time.time() - start_time
                        
                        # Complete progress
                        progress_bar.progress(100)
                        agent_status_placeholder.empty()  # Clear the agent status
                        
                        # Save results to session state
                        st.session_state.teaching_plan_execution_time = execution_time
                        st.session_state.teaching_plan_result = result_text
                        st.session_state.teaching_plan_generated = True
                        
                        # Display the result
                        st.success(f"✅ Teaching plan generated successfully in {execution_time:.2f} seconds!")
                        st.subheader("📝 Generated Teaching Plan")
                        st.markdown(result_text)
                        
                        # Generate appropriate filename based on content type
                        if st.session_state.content_type["type"] == "Case Study":
                            filename = "case_study_teaching_plan.md"
                        elif st.session_state.content_type["type"] == "Scientific Article":
                            filename = "article_teaching_plan.md"
                        elif st.session_state.content_type["type"] == "News Item":
                            filename = "news_discussion_guide.md"
                        else:
                            filename = "teaching_plan.md"
                        
                        # Provide download option
                        st.download_button(
                            label="📥 Download Teaching Plan (Markdown)",
                            data=result_text,
                            file_name=filename,
                            mime="text/markdown",
                        )
                        
                    except Exception as e:
                        st.error(f"❌ An error occurred during processing: {str(e)}")
                
                # Show instructions
                with st.expander("ℹ️ How it works"):
                    if st.session_state.content_type["type"] == "Case Study":
                        instructions = """
                        The Teaching Plan Generator creates a comprehensive 2-hour lesson plan for case studies with:
                        
                        1. **Clear learning objectives** tied to the case study
                        2. **Structured timeline** with time allocations for each activity
                        3. **Case discussion framework** for effective analysis
                        4. **Group activities** that apply business concepts
                        5. **Assessment strategies** to measure understanding
                        """
                    elif st.session_state.content_type["type"] == "Scientific Article":
                        instructions = """
                        The Teaching Plan Generator creates a comprehensive 2-hour lesson plan for scientific articles with:
                        
                        1. **Research objectives** tied to the article's methodology
                        2. **Structured timeline** for analyzing research components
                        3. **Critical evaluation framework** for assessing research quality
                        4. **Application activities** that connect research to practice
                        5. **Research extension opportunities** for further exploration
                        """
                    elif st.session_state.content_type["type"] == "News Item":
                        instructions = """
                        The Discussion Guide Generator creates a comprehensive 2-hour lesson plan for news items with:
                        
                        1. **Context objectives** for understanding the news event
                        2. **Structured timeline** for exploring different perspectives
                        3. **Media literacy framework** for critical news consumption
                        4. **Current events activities** that connect to broader concepts
                        5. **Discussion strategies** to explore implications and impacts
                        """
                    else:
                        instructions = """
                        The Teaching Guide Generator creates a comprehensive 2-hour lesson plan with:
                        
                        1. **Clear learning objectives** tied to the document content
                        2. **Structured timeline** with time allocations for each activity
                        3. **Engaging activities** for effective student learning
                        4. **Discussion questions** to promote critical thinking
                        5. **Assessment strategies** to measure understanding
                        """
                    
                    st.write(instructions)
                    
                    st.write("""
                    The AI uses a team of specialized agents to analyze your document, create a draft plan, 
                    review it for quality, and finalize it into a polished teaching resource.
                    """)
        else:
            # Display the previously generated content
            st.success(f"✅ Teaching plan generated successfully in {st.session_state.teaching_plan_execution_time:.2f} seconds!")
            st.subheader("📝 Generated Teaching Plan")
            st.markdown(st.session_state.teaching_plan_result)
            
            # Generate appropriate filename based on content type
            # Generate appropriate filename based on content type
            if st.session_state.content_type["type"] == "Case Study":
                md_filename = "case_study_teaching_plan.md"
                docx_filename = "case_study_teaching_plan.docx"
                docx_title = "Case Study Teaching Plan"
            elif st.session_state.content_type["type"] == "Scientific Article":
                md_filename = "article_teaching_plan.md"
                docx_filename = "article_teaching_plan.docx"
                docx_title = "Scientific Article Teaching Plan"
            elif st.session_state.content_type["type"] == "News Item":
                md_filename = "news_discussion_guide.md"
                docx_filename = "news_discussion_guide.docx"
                docx_title = "News Item Discussion Guide"
            else:
                md_filename = "teaching_plan.md"
                docx_filename = "teaching_plan.docx"
                docx_title = "Teaching Plan"
            
            # Create columns for download buttons
            col1, col2 = st.columns(2)
            
            # Provide download options
            with col1:
                st.download_button(
                    label="📥 Download as Markdown",
                    data=st.session_state.teaching_plan_result,
                    file_name=md_filename,
                    mime="text/markdown",
                )
            
            with col2:
                # Create DOCX file from the markdown content
                docx_path = create_docx_from_markdown(
                    st.session_state.teaching_plan_result,
                    docx_title,
                    st.session_state.content_type["type"]
                )
                
                # Provide DOCX download button
                with open(docx_path, 'rb') as file:
                    st.download_button(
                        label="📥 Download as DOCX",
                        data=file,
                        file_name=docx_filename,
                        mime="application/vnd.openxmlformats-officedocument.wordprocessingml.document",
                        key="teaching_plan_docx"
                    )
    
    #------------------ Discussion Framework Generator Tab ------------------#
    with tab3:
        # Update header based on content type
        st.header(tab_names[2])
        
        # Update description based on content type
        if st.session_state.content_type["type"] == "Case Study":
            description = "Generate a structured board plan analysis for the case study."
        elif st.session_state.content_type["type"] == "Scientific Article":
            description = "Generate key discussion points based on the scientific article."
        elif st.session_state.content_type["type"] == "News Item":
            description = "Generate a current events analysis framework for the news item."
        else:
            description = "Generate a structured discussion framework for the document."
            
        st.write(description)
        
        if not st.session_state.board_plan_generated:
            if not api_key:
                st.warning("⚠️ Please enter an API key in the sidebar before proceeding.")
            else:
                # Create a button to start generation with label based on content type
                if st.button(f"Generate {tab_names[2]}", key="board_plan_button"):
                    try:
                        # Initialize the analyzer
                        analyzer = DiscussionFrameworkAnalyzer(api_key)
                        
                        with st.spinner(f"Analyzing {st.session_state.content_type['type'].lower()} to generate {tab_names[2].lower()}..."):
                            # Get content from first uploaded file
                            file = st.session_state.uploaded_files[0]
                            file_content = file.getvalue()
                            
                            # Process the document with content type awareness
                            analysis_result = analyzer.process_document(
                                file_content, 
                                st.session_state.content_type["type"]
                            )
                            
                            # Store in session state
                            st.session_state.board_plan_result = analysis_result
                            st.session_state.board_plan_generated = True
                            
                            # Generate markdown for download
                            markdown_content = f"# {tab_names[2]}\n\n"
                            for section in analysis_result['sections']:
                                markdown_content += f"## {section['title']}\n\n"
                                for point in section['points']:
                                    markdown_content += f"* {point}\n"
                                markdown_content += "\n"
                            
                            if 'questions' in analysis_result:
                                markdown_content += f"## {analysis_result['questions']['title']}\n\n"
                                for point in analysis_result['questions']['points']:
                                    markdown_content += f"* {point}\n"
                            
                            st.session_state.board_plan_markdown = markdown_content
                        
                        st.success(f"{tab_names[2]} generated successfully!")
                        
                        # Display the analysis results
                        st.subheader(tab_names[2])
                        for section in analysis_result['sections']:
                            with st.expander(section['title'], expanded=True):
                                for point in section['points']:
                                    st.markdown(f"• {point}")
                        
                        if 'questions' in analysis_result:
                            st.subheader("Questions for Further Exploration")
                            for point in analysis_result['questions']['points']:
                                st.markdown(f"• {point}")

                        if 'discussion' in analysis_result:
                            st.subheader("Simulated Discussion")
                            st.markdown(analysis_result['discussion']['content'])

                        # Generate appropriate filename based on content type
                        if st.session_state.content_type["type"] == "Case Study":
                            md_filename = "board_plan.md"
                            docx_filename = "board_plan.docx"
                            docx_title = "Board Plan Analysis"
                        elif st.session_state.content_type["type"] == "Scientific Article":
                            md_filename = "discussion_points.md"
                            docx_filename = "discussion_points.docx"
                            docx_title = "Discussion Points Analysis"
                        elif st.session_state.content_type["type"] == "News Item":
                            md_filename = "current_events_framework.md"
                            docx_filename = "current_events_framework.docx"
                            docx_title = "Current Events Framework"
                        else:
                            md_filename = "discussion_framework.md"
                            docx_filename = "discussion_framework.docx"
                            docx_title = "Discussion Framework"
                        
                        # Create columns for download buttons
                        col1, col2 = st.columns(2)
                        
                        # Provide download options
                        with col1:
                            st.download_button(
                                label="📥 Download as Markdown",
                                data=st.session_state.board_plan_markdown,
                                file_name=md_filename,
                                mime="text/markdown",
                            )
                        
                        with col2:
                            # Create DOCX file from the markdown content
                            docx_path = create_docx_from_markdown(
                                st.session_state.board_plan_markdown,
                                docx_title,
                                st.session_state.content_type["type"]
                            )
                            
                            # Provide DOCX download button
                            with open(docx_path, 'rb') as file:
                                st.download_button(
                                    label="📥 Download as DOCX",
                                    data=file,
                                    file_name=docx_filename,
                                    mime="application/vnd.openxmlformats-officedocument.wordprocessingml.document",
                                    key="discussion_framework_docx"
                                )
                        
                    except Exception as e:
                        st.error(f"An error occurred: {str(e)}")
                        st.error("Please ensure the file is not encrypted and contains extractable text.")
                
                # Show instructions based on content type
                with st.expander("ℹ️ How it works"):
                    if st.session_state.content_type["type"] == "Case Study":
                        instructions = """
                        The Board Plan Generator creates a structured analysis with these components:
                        
                        1. **Main business challenge/opportunity** being addressed in the case study
                        2. **Industry context and competitors** relevant to the case
                        3. **Key decision points** that require strategic choices
                        4. **Stakeholder perspectives** from different roles and viewpoints
                        5. **Implementation considerations** including resources, timeline, and metrics
                        6. **Alternative approaches** that could be considered
                        7. **Learning takeaways** for applying key lessons
                        """
                    elif st.session_state.content_type["type"] == "Scientific Article":
                        instructions = """
                        The Discussion Points Generator creates a structured analysis with these components:
                        
                        1. **Main research question/hypothesis** examined in the article
                        2. **Methodology assessment** evaluating research approach
                        3. **Key findings summary** highlighting important discoveries
                        4. **Limitations and validity** considerations
                        5. **Theoretical implications** for the field
                        6. **Practical applications** of the research
                        7. **Future research directions** suggested by the findings
                        8. **Interdisciplinary connections** to other fields
                        """
                    elif st.session_state.content_type["type"] == "News Item":
                        instructions = """
                        The Current Events Framework Generator creates a structured analysis with these components:
                        
                        1. **Main event/development** described in the news item
                        2. **Historical and political context** surrounding the event
                        3. **Key stakeholders and perspectives** involved
                        4. **Media coverage analysis** examining reporting approaches
                        5. **Immediate impacts** on various sectors
                        6. **Long-term implications** for policy and society
                        7. **Related developments** connected to this news
                        8. **Critical thinking questions** for deeper analysis
                        """
                    else:
                        instructions = """
                        The Discussion Framework Generator creates a structured analysis with these components:
                        
                        1. **Main topic/theme** of the document
                        2. **Key concepts** central to understanding the content
                        3. **Important relationships** between ideas and elements
                        4. **Supporting evidence** presented in the document
                        5. **Counterarguments/limitations** worth considering
                        6. **Practical applications** of the content
                        7. **Discussion questions** for deeper exploration
                        8. **Further exploration topics** related to the content
                        """
                    
                    st.write(instructions)
                    
                    st.write("""
                    The generator provides a structured framework perfect for teaching, presentations, 
                    or educational discussions.
                    """)
        else:
            # Display the previously generated content
            st.success(f"{tab_names[2]} generated successfully!")
            
            # Display the analysis results
            st.subheader(tab_names[2])
            for section in st.session_state.board_plan_result['sections']:
                with st.expander(section['title'], expanded=True):
                    for point in section['points']:
                        st.markdown(f"• {point}")
            
            if 'questions' in st.session_state.board_plan_result:
                st.subheader("Questions for Further Exploration")
                for point in st.session_state.board_plan_result['questions']['points']:
                    st.markdown(f"• {point}")
            
            # Generate appropriate filename based on content type
            if st.session_state.content_type["type"] == "Case Study":
                filename = "board_plan.md"
            elif st.session_state.content_type["type"] == "Scientific Article":
                filename = "discussion_points.md"
            elif st.session_state.content_type["type"] == "News Item":
                filename = "current_events_framework.md"
            else:
                filename = "discussion_framework.md"
                
            # Provide download option
            st.download_button(
                label=f"📥 Download {tab_names[2]} (Markdown)",
                data=st.session_state.board_plan_markdown,
                file_name=filename,
                mime="text/markdown",
            )

else:
    # Display welcome message and instructions when no files are uploaded
    st.markdown("""
    ## Welcome to the Document Analysis Suite
    
    This application provides three powerful tools for analyzing different types of documents:
    
    1. **Document Analysis Generator**: Creates a comprehensive analysis with sections tailored to the document type (case study, scientific article, news item, or other content).
    
    2. **Teaching Plan Generator**: Develops a 2-hour lesson plan with activities, discussions, and assessments based on the document.
    
    3. **Discussion Framework Generator**: Produces a structured analysis with key insights based on the type of document uploaded.
    
    ### Getting Started
    
    1. Enter your API key in the sidebar
    2. Upload one or more document files in PDF, DOCX, PPT, or PPTX format
    3. The system will automatically detect the document type (case study, scientific article, news item, or other)
    4. You can confirm or change the detected document type
    5. Select the generator tab you want to use
    6. Click the generate button and wait for the results
    
    Each generator produces downloadable content tailored to your document type.
    """)

# Footer
st.divider()
st.caption("Created with CrewAI, Streamlit, and Gemini • Built by Arun Kashyap • © 2025")               
